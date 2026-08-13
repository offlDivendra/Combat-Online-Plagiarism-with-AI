"""
PowerPoint Generator for AI Plagiarism Detection System
Uses python-pptx to create a professional 20-slide presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette ────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x1A, 0x6E)   # deep navy (title bg)
MID_BLUE    = RGBColor(0x1E, 0x3A, 0x8A)   # section headers
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # accent / bullets
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF1, 0xF5, 0xF9)
DARK_GRAY   = RGBColor(0x1E, 0x29, 0x3B)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)
RED         = RGBColor(0xDC, 0x26, 0x26)
YELLOW      = RGBColor(0xCA, 0x8A, 0x04)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper utilities ──────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color and line_width:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=16, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None):
    """Dark blue top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.35, fill_color=DARK_BLUE)
    add_text(slide, title, 0.4, 0.1, 12, 0.9,
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.88, 12, 0.45,
                 font_size=15, bold=False, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
    # thin accent line
    add_rect(slide, 0, 1.35, 13.33, 0.05, fill_color=ACCENT_BLUE)


def content_bg(slide):
    """Light background for content area."""
    add_rect(slide, 0, 1.4, 13.33, 6.1, fill_color=LIGHT_GRAY)


def bullet_box(slide, items, l, t, w, h,
               font_size=16, title=None, title_color=MID_BLUE,
               box_color=WHITE, bullet="●"):
    """White card with optional title and bullet list."""
    add_rect(slide, l, t, w, h, fill_color=box_color,
             line_color=ACCENT_BLUE, line_width=0.75)
    txBox = slide.shapes.add_textbox(
        Inches(l + 0.15), Inches(t + 0.1),
        Inches(w - 0.3), Inches(h - 0.2))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        add_para(tf, title, font_size=font_size + 2,
                 bold=True, color=title_color, space_before=0)
    for item in items:
        add_para(tf, f"{bullet}  {item}", font_size=font_size,
                 color=DARK_GRAY, space_before=4)
    return txBox


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=DARK_BLUE)
# diagonal accent strip
add_rect(s, 0, 4.5, 13.33, 0.08, fill_color=ACCENT_BLUE)
add_rect(s, 0, 4.6, 13.33, 0.04, fill_color=RGBColor(0x93, 0xC5, 0xFD))

add_text(s, "COMBAT ONLINE PLAGIARISM", 0.6, 1.0, 12.0, 1.0,
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "WITH ARTIFICIAL INTELLIGENCE", 0.6, 1.9, 12.0, 0.9,
         font_size=36, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "An AI-Powered Plagiarism Detection System", 0.6, 2.9, 12.0, 0.6,
         font_size=20, bold=False, color=RGBColor(0xBA, 0xD4, 0xFB),
         align=PP_ALIGN.CENTER)

add_rect(s, 3.5, 3.7, 6.33, 0.05, fill_color=ACCENT_BLUE)

add_text(s, "BTech CSE Final Year Project  ●  2026", 0.6, 3.9, 12.0, 0.5,
         font_size=16, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
add_text(s, "Technology Stack:  Python  ·  FastAPI  ·  React.js  ·  NLP  ·  SQLite",
         0.6, 4.5, 12.0, 0.5, font_size=14,
         color=RGBColor(0x7D, 0xD3, 0xFC), align=PP_ALIGN.CENTER)

# Bottom strip
add_rect(s, 0, 6.8, 13.33, 0.7, fill_color=RGBColor(0x0F, 0x0F, 0x4A))
add_text(s, "TF-IDF  ●  N-Gram  ●  Fuzzy Matching  ●  PDF Reports  ●  REST API",
         0.6, 6.85, 12.0, 0.5, font_size=13,
         color=RGBColor(0x7D, 0xD3, 0xFC), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Table of Contents
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Table of Contents", "What we will cover today")

topics_left = [
    "01  Project Overview & Problem Statement",
    "02  Objectives & Scope",
    "03  Technologies Used",
    "04  System Architecture",
    "05  NLP Pipeline Explained",
    "06  Algorithm 1 — TF-IDF",
    "07  Algorithm 2 — N-Gram Similarity",
    "08  Algorithm 3 — Fuzzy Matching",
    "09  Combined Scoring Engine",
    "10  Database Design",
]
topics_right = [
    "11  API Endpoints",
    "12  Frontend Pages",
    "13  Project Folder Structure",
    "14  How to Run the Project",
    "15  Test Results & Evaluation",
    "16  Live Demo Walkthrough",
    "17  Limitations",
    "18  Future Enhancements",
    "19  Viva Q&A Highlights",
    "20  Conclusion",
]

bullet_box(s, topics_left, 0.3, 1.55, 6.2, 5.7,
           font_size=14, bullet="▸", box_color=WHITE)
bullet_box(s, topics_right, 6.8, 1.55, 6.2, 5.7,
           font_size=14, bullet="▸", box_color=WHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — Project Overview & Problem Statement
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Project Overview", "What is this project and why does it matter?")

add_rect(s, 0.3, 1.55, 8.0, 1.6, fill_color=MID_BLUE)
add_text(s, "What Is This Project?", 0.5, 1.6, 7.6, 0.5,
         font_size=17, bold=True, color=WHITE)
add_text(s,
    "An AI-powered web application that automatically detects plagiarism "
    "in documents by comparing text using three NLP algorithms and returning "
    "a combined similarity score with detailed evidence.",
    0.5, 2.05, 7.6, 1.0, font_size=14, color=WHITE, wrap=True)

bullet_box(s,
    ["Manual checking takes hours for teachers",
     "Students paraphrase to evade simple tools",
     "No documentation / evidence trail exists",
     "No centralized system across departments"],
    0.3, 3.35, 5.8, 3.8,
    font_size=14, title="❌  The Problem", title_color=RED)

bullet_box(s,
    ["Check documents in seconds, not hours",
     "Detect paraphrasing with TF-IDF",
     "Generate PDF evidence reports",
     "Store full analysis history in DB"],
    6.3, 3.35, 6.7, 3.8,
    font_size=14, title="✅  Our Solution", title_color=GREEN)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — Objectives & Scope
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Objectives & Scope", "What we set out to achieve")

objectives = [
    "Implement multi-algorithm NLP pipeline (TF-IDF + N-gram + Fuzzy)",
    "Build a RESTful API backend using FastAPI (Python)",
    "Create an interactive React.js frontend with 5 pages",
    "Support multiple file formats: TXT, PDF, DOCX",
    "Generate professional PDF plagiarism reports",
    "Store analysis history in SQLite database",
    "Achieve >85% accuracy on test document set",
    "Write 100+ unit tests for code reliability",
]

scope_in = [
    "Text-based document comparison",
    "English language documents",
    "Local reference document database",
    "Web-based interface",
    "PDF report generation",
]

scope_out = [
    "Internet / web crawling",
    "Code plagiarism detection",
    "Multi-language support",
    "User authentication system",
    "Mobile application",
]

bullet_box(s, objectives, 0.3, 1.55, 7.8, 5.7,
           font_size=13, title="🎯  Project Objectives", title_color=MID_BLUE)
bullet_box(s, scope_in,   8.3, 1.55, 4.7, 3.1,
           font_size=13, title="✅  In Scope", title_color=GREEN)
bullet_box(s, scope_out,  8.3, 4.85, 4.7, 2.4,
           font_size=13, title="❌  Out of Scope", title_color=RED)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Technologies Used
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Technologies Used", "The tools and libraries powering the system")

tech_data = [
    ("Python 3.14",      "Core programming language",         MID_BLUE),
    ("FastAPI",          "REST API framework — fast & modern", MID_BLUE),
    ("React.js + Vite",  "Frontend UI framework",              RGBColor(0x06,0x7B,0xA5)),
    ("scikit-learn",     "TF-IDF vectorization & cosine sim",  RGBColor(0x0E,0x73,0x48)),
    ("NLTK",             "Tokenization & stopword removal",    RGBColor(0x0E,0x73,0x48)),
    ("RapidFuzz",        "Ultra-fast fuzzy string matching",   RGBColor(0x0E,0x73,0x48)),
    ("SQLAlchemy",       "ORM — Python ↔ SQLite database",    RGBColor(0x71,0x3F,0x12)),
    ("ReportLab",        "PDF report generation",              RGBColor(0x71,0x3F,0x12)),
    ("Axios",            "HTTP requests from React to API",    RGBColor(0x06,0x7B,0xA5)),
    ("pytest",           "100+ unit tests",                    RGBColor(0x4C,0x0D,0x99)),
]

cols = 2
rows_per_col = 5
box_w, box_h = 5.9, 0.9
gap_x, gap_y = 0.25, 0.1
start_x = [0.3, 6.7]

for i, (name, desc, color) in enumerate(tech_data):
    col = i // rows_per_col
    row = i %  rows_per_col
    x = start_x[col]
    y = 1.6 + row * (box_h + gap_y)
    add_rect(s, x, y, box_w, box_h, fill_color=color)
    add_text(s, name, x + 0.12, y + 0.04, 2.2, 0.45,
             font_size=15, bold=True, color=WHITE)
    add_text(s, desc, x + 0.12, y + 0.42, box_w - 0.25, 0.4,
             font_size=12, color=RGBColor(0xD1,0xFA,0xFF), wrap=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "System Architecture", "3-tier architecture: Frontend → Backend → Database")

layers = [
    (ACCENT_BLUE,  "LAYER 1 — FRONTEND  (React.js + Vite  ·  Port 5173)",
     "Home  ·  Analyze  ·  Results  ·  History  ·  Documents  |  Axios HTTP calls"),
    (MID_BLUE,     "LAYER 2 — BACKEND   (FastAPI + Python  ·  Port 8000)",
     "Routes: /analyze  /documents  /reports  |  Services: TF-IDF · N-gram · Fuzzy · PDF"),
    (DARK_BLUE,    "LAYER 3 — DATABASE  (SQLite + SQLAlchemy)",
     "Tables: reference_documents  ·  analysis_history  ·  analysis_reports"),
]

for i, (color, title, detail) in enumerate(layers):
    y = 1.65 + i * 1.55
    add_rect(s, 0.5, y, 12.3, 1.3, fill_color=color)
    add_text(s, title,  0.7, y + 0.06, 11.8, 0.55, font_size=16, bold=True,  color=WHITE)
    add_text(s, detail, 0.7, y + 0.58, 11.8, 0.55, font_size=13, bold=False,
             color=RGBColor(0xBA,0xD4,0xFB), wrap=True)
    if i < 2:
        add_text(s, "▼  HTTP / ORM", 5.7, y + 1.28, 2.5, 0.3,
                 font_size=13, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

flow_items = [
    "1. User submits text or file",
    "2. React sends POST /api/analyze",
    "3. FastAPI loads reference docs from SQLite",
    "4. Runs 3 NLP algorithms in sequence",
    "5. Calculates weighted combined score",
    "6. Saves result → returns JSON",
    "7. React renders Results page",
]
bullet_box(s, flow_items, 0.5, 6.15, 12.3, 1.1,
           font_size=12, title=None, box_color=WHITE, bullet="→")

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — NLP Pipeline
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "NLP Pipeline", "How raw text becomes a similarity score — step by step")

steps = [
    (ACCENT_BLUE, "STEP 1", "RAW INPUT",
     '"Machine Learning is a branch of AI that helps computers LEARN!"'),
    (MID_BLUE,    "STEP 2", "NORMALIZE",
     '"machine learning is a branch of ai that helps computers learn"'),
    (RGBColor(0x0E,0x73,0x48), "STEP 3", "TOKENIZE",
     '["machine", "learning", "is", "a", "branch", "of", "ai", "helps", "computers", "learn"]'),
    (RGBColor(0x71,0x3F,0x12), "STEP 4", "REMOVE STOPWORDS",
     '["machine", "learning", "branch", "ai", "helps", "computers", "learn"]'),
    (RGBColor(0x4C,0x0D,0x99), "STEP 5", "ALGORITHMS",
     "TF-IDF  +  N-gram  +  Fuzzy  →  Combined Score  →  Classification"),
]

for i, (color, step, title, content) in enumerate(steps):
    y = 1.6 + i * 1.05
    add_rect(s, 0.3, y, 1.15, 0.88, fill_color=color)
    add_text(s, step,  0.3,  y + 0.04, 1.15, 0.35, font_size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, 0.3,  y + 0.42, 1.15, 0.38, font_size=10,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 1.6, y, 11.3, 0.88, fill_color=WHITE,
             line_color=color, line_width=1)
    add_text(s, content, 1.75, y + 0.1, 11.0, 0.72,
             font_size=13, color=DARK_GRAY, wrap=True)
    if i < 4:
        add_text(s, "▼", 0.7, y + 0.9, 0.6, 0.2,
                 font_size=14, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — TF-IDF Algorithm
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Algorithm 1 — TF-IDF Similarity",
             "Term Frequency × Inverse Document Frequency  ·  Weight: 50%")

add_rect(s, 0.3, 1.6, 12.7, 0.65, fill_color=MID_BLUE)
add_text(s, "What is TF-IDF?", 0.5, 1.65, 12.0, 0.5,
         font_size=17, bold=True, color=WHITE)

bullet_box(s,
    ["TF (Term Frequency) = How often a word appears in the document",
     "IDF (Inverse Doc Freq) = How RARE the word is across ALL documents",
     "TF-IDF Score = TF × IDF   →   Rare but frequent words score highest",
     "Example: 'blockchain' scores high; 'the' scores near zero"],
    0.3, 2.35, 8.2, 2.6, font_size=13, title="How It Works", title_color=MID_BLUE)

bullet_box(s,
    ["Convert text to TF-IDF number vector",
     "cos(θ) = (A · B) / (||A|| × ||B||)",
     "Score 0 = completely different",
     "Score 1 = identical documents"],
    8.7, 2.35, 4.3, 2.6, font_size=13, title="Cosine Similarity", title_color=MID_BLUE)

add_rect(s, 0.3, 5.15, 12.7, 0.5, fill_color=RGBColor(0xDB,0xEA,0xFE))
add_text(s, "Why 50% weight?  →  Best at capturing semantic / topic-level similarity."
            "  Even paraphrased content scores well if same vocabulary is used.",
         0.5, 5.18, 12.2, 0.42, font_size=13, color=MID_BLUE, wrap=True)

examples = [
    ("Text A", "Machine learning is a subset of AI.", "Text B", "ML is part of artificial intelligence.", "82%", GREEN),
    ("Text A", "Rome was an ancient civilization.", "Text B", "Blockchain uses distributed ledgers.", "4%", RED),
]
for i, (l1, t1, l2, t2, score, col) in enumerate(examples):
    x = 0.3 + i * 6.4
    add_rect(s, x, 5.8, 6.1, 1.45, fill_color=WHITE, line_color=ACCENT_BLUE, line_width=0.75)
    add_text(s, f"{l1}: {t1}", x+0.1, 5.88, 5.0, 0.4, font_size=12, color=DARK_GRAY, wrap=True)
    add_text(s, f"{l2}: {t2}", x+0.1, 6.28, 5.0, 0.4, font_size=12, color=DARK_GRAY, wrap=True)
    add_text(s, f"Similarity: {score}", x+0.1, 6.7, 3.0, 0.4, font_size=14, bold=True, color=col)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — N-Gram Algorithm
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Algorithm 2 — N-Gram Similarity",
             "Detecting copied phrases using word sequences  ·  Weight: 20%")

add_rect(s, 0.3, 1.6, 12.7, 0.65, fill_color=MID_BLUE)
add_text(s, "What is an N-Gram?  →  A sequence of N consecutive words in a text.",
         0.5, 1.65, 12.2, 0.5, font_size=16, bold=True, color=WHITE)

gram_data = [
    ("Unigrams\n(N=1)", '["artificial", "intelligence",\n"is", "amazing"]', ACCENT_BLUE),
    ("Bigrams\n(N=2)", '["artificial intelligence",\n"intelligence is", "is amazing"]', MID_BLUE),
    ("Trigrams\n(N=3)", '["artificial intelligence is",\n"intelligence is amazing"]', DARK_BLUE),
]
for i, (label, content, color) in enumerate(gram_data):
    x = 0.3 + i * 4.3
    add_rect(s, x, 2.45, 4.0, 1.8, fill_color=color)
    add_text(s, label,   x+0.1, 2.52, 3.7, 0.65, font_size=15, bold=True, color=WHITE)
    add_text(s, content, x+0.1, 3.1,  3.7, 1.0,  font_size=12, color=RGBColor(0xD1,0xFA,0xFF), wrap=True)

add_rect(s, 0.3, 4.45, 12.7, 0.55, fill_color=RGBColor(0xDB,0xEA,0xFE))
add_text(s, "We use Bigrams (N=2)  ·  Jaccard Similarity = Common N-grams ÷ Total unique N-grams",
         0.5, 4.48, 12.2, 0.45, font_size=14, bold=True, color=MID_BLUE)

bullet_box(s,
    ['Text 1 bigrams: {"machine learning", "learning is", "is fast"}',
     'Text 2 bigrams: {"machine learning", "learning can", "can fail"}',
     'Common bigrams: {"machine learning"}  =  1',
     'Total unique:   5  →  Jaccard = 1/5 = 20%'],
    0.3, 5.1, 7.8, 2.1, font_size=13, title="Worked Example", title_color=MID_BLUE)

bullet_box(s,
    ["Catches exact phrase copying",
     "Low score for paraphrased text",
     "Fast to compute",
     "Works with any language"],
    8.3, 5.1, 4.7, 2.1, font_size=13, title="Pros & Cons", title_color=MID_BLUE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Fuzzy Matching Algorithm
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Algorithm 3 — Fuzzy Matching",
             "Sentence-level similarity with tolerance for variations  ·  Weight: 30%")

add_rect(s, 0.3, 1.6, 12.7, 0.65, fill_color=MID_BLUE)
add_text(s,
    "Fuzzy matching compares strings character-by-character, "
    "tolerating typos, word order changes, and minor modifications.",
    0.5, 1.65, 12.2, 0.5, font_size=15, bold=True, color=WHITE, wrap=True)

methods = [
    ("Simple Ratio",    "Direct character comparison", "90%",  ACCENT_BLUE),
    ("Token Sort",      "Compare after sorting words", "95%",  MID_BLUE),
    ("Token Set",       "Common words ÷ all words",   "88%",  RGBColor(0x0E,0x73,0x48)),
    ("WRatio",          "Smart weighted combination",  "93%",  RGBColor(0x71,0x3F,0x12)),
]
for i, (name, desc, score, color) in enumerate(methods):
    x = 0.3 + i * 3.2
    add_rect(s, x, 2.45, 3.0, 1.4, fill_color=color)
    add_text(s, name,  x+0.1, 2.52, 2.8, 0.45, font_size=14, bold=True, color=WHITE)
    add_text(s, desc,  x+0.1, 2.92, 2.8, 0.55, font_size=11, color=RGBColor(0xD1,0xFA,0xFF), wrap=True)
    add_text(s, score, x+0.1, 3.48, 2.8, 0.3,  font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

bullet_box(s,
    ["Split both documents into individual sentences",
     "Compare each query sentence vs all reference sentences",
     "Keep matches that score above 82% threshold",
     "Final Score = (matched sentences / total sentences) × avg match quality",
     "Example: 8/10 matched at 90% avg → 8/10 × 90 = 72%"],
    0.3, 4.05, 8.5, 3.15, font_size=13, title="How Sentence Coverage Works", title_color=MID_BLUE)

bullet_box(s,
    ["Catches minor modifications",
     "Handles word order changes",
     "Detects typo-masked plagiarism",
     "Threshold: 82% to avoid false positives"],
    9.0, 4.05, 4.0, 3.15, font_size=13, title="Why Useful?", title_color=GREEN)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — Combined Scoring Engine
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Combined Scoring Engine", "How three scores become one final verdict")

algo_data = [
    ("TF-IDF",   "50%", "89.87",  "44.94",  ACCENT_BLUE),
    ("N-Gram",   "20%", "64.13",  "12.83",  MID_BLUE),
    ("Fuzzy",    "30%", "99.37",  "29.81",  RGBColor(0x0E,0x73,0x48)),
]
headers = ["Algorithm", "Weight", "Raw Score", "Weighted Score"]
col_x   = [0.4, 3.2, 5.8, 8.6]
col_w   = [2.6, 2.4, 2.6, 2.8]

add_rect(s, 0.3, 1.65, 11.8, 0.55, fill_color=DARK_BLUE)
for j, h in enumerate(headers):
    add_text(s, h, col_x[j], 1.68, col_w[j], 0.45,
             font_size=14, bold=True, color=WHITE)

for i, (name, weight, raw, weighted, color) in enumerate(algo_data):
    y = 2.3 + i * 0.72
    bg = WHITE if i % 2 == 0 else RGBColor(0xF8,0xFA,0xFF)
    add_rect(s, 0.3, y, 11.8, 0.65, fill_color=bg,
             line_color=ACCENT_BLUE, line_width=0.5)
    add_rect(s, 0.3, y, 0.12, 0.65, fill_color=color)
    vals = [name, weight, f"{raw}%", f"{weighted}"]
    for j, v in enumerate(vals):
        clr = color if j == 0 else DARK_GRAY
        add_text(s, v, col_x[j], y + 0.12, col_w[j], 0.42,
                 font_size=14, bold=(j == 0), color=clr)

add_rect(s, 0.3, 4.5, 11.8, 0.62, fill_color=DARK_BLUE)
add_text(s, "COMBINED SCORE  =  44.94 + 12.83 + 29.81  =  87.57%   →   Potential Plagiarism  🔴",
         0.5, 4.55, 11.4, 0.5, font_size=16, bold=True, color=WHITE)

thresholds = [
    ("0 – 20%",   "Mostly Original",        GREEN),
    ("20 – 40%",  "Low Similarity",          RGBColor(0x84,0xCC,0x16)),
    ("40 – 60%",  "Moderate Similarity",     YELLOW),
    ("60 – 80%",  "High Similarity",         ORANGE),
    ("80 – 100%", "Potential Plagiarism",    RED),
]
add_text(s, "Classification Thresholds:", 0.4, 5.3, 6.0, 0.4,
         font_size=15, bold=True, color=MID_BLUE)
for i, (rng, label, color) in enumerate(thresholds):
    x = 0.3 + i * 2.55
    add_rect(s, x, 5.75, 2.35, 0.9, fill_color=color)
    add_text(s, rng,   x+0.08, 5.8,  2.2, 0.38, font_size=13, bold=True,  color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, label, x+0.08, 6.18, 2.2, 0.42, font_size=11, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — Database Design
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Database Design", "SQLite + SQLAlchemy ORM — 3 tables")

tables = [
    ("reference_documents",
     ["id (PK)", "filename", "original_filename", "file_extension",
      "file_size", "text_content", "word_count", "sentence_count", "upload_date"],
     MID_BLUE),
    ("analysis_history",
     ["id (PK)", "document_name", "analysis_date", "overall_similarity",
      "tfidf_score", "ngram_score", "fuzzy_score", "classification",
      "total_matches", "sources (JSON)", "sentence_matches (JSON)", "submitted_text"],
     RGBColor(0x0E,0x73,0x48)),
    ("analysis_reports",
     ["id (PK)", "analysis_id (FK→history)", "report_filename",
      "generation_date", "report_format", "file_path", "file_size"],
     RGBColor(0x71,0x3F,0x12)),
]

start_x = [0.25, 4.6, 8.95]
for i, (name, cols, color) in enumerate(tables):
    x = start_x[i]
    add_rect(s, x, 1.65, 4.1, 0.55, fill_color=color)
    add_text(s, name, x+0.1, 1.68, 3.9, 0.45, font_size=14, bold=True, color=WHITE)
    for j, col in enumerate(cols):
        bg = WHITE if j % 2 == 0 else LIGHT_GRAY
        y = 2.25 + j * 0.48
        add_rect(s, x, y, 4.1, 0.46, fill_color=bg,
                 line_color=color, line_width=0.4)
        icon = "🔑 " if "PK" in col else ("🔗 " if "FK" in col else "   ")
        add_text(s, icon + col, x+0.1, y+0.06, 3.9, 0.35,
                 font_size=12, color=DARK_GRAY)

add_rect(s, 0.25, 6.8, 12.8, 0.45, fill_color=RGBColor(0xDB,0xEA,0xFE))
add_text(s, "relationship:  analysis_history  ──has many──►  analysis_reports  "
            "  |  reference_documents  ──used by──►  analysis_history",
         0.4, 6.83, 12.4, 0.38, font_size=12, color=MID_BLUE, wrap=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — API Endpoints
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "API Endpoints", "RESTful API built with FastAPI  ·  Auto-docs at /api/docs")

endpoints = [
    ("POST",   "/api/analyze",              "Analyze pasted text for plagiarism",     GREEN),
    ("POST",   "/api/analyze/file",         "Analyze uploaded TXT / PDF / DOCX",      GREEN),
    ("GET",    "/api/analyze/history",      "Retrieve all past analysis records",     ACCENT_BLUE),
    ("GET",    "/api/analyze/history/{id}", "Get detailed result for one analysis",   ACCENT_BLUE),
    ("DELETE", "/api/analyze/history/{id}", "Delete an analysis record",              RED),
    ("GET",    "/api/analyze/statistics",   "Get system-wide stats & counts",         ACCENT_BLUE),
    ("GET",    "/api/documents",            "List all uploaded reference documents",  MID_BLUE),
    ("POST",   "/api/documents/upload",     "Upload a new reference document",        GREEN),
    ("DELETE", "/api/documents/{id}",       "Remove a reference document",            RED),
    ("POST",   "/api/reports/generate",     "Generate a PDF report for an analysis",  RGBColor(0x71,0x3F,0x12)),
    ("GET",    "/api/reports/download/{id}","Download the generated PDF report",      RGBColor(0x71,0x3F,0x12)),
]
method_colors = {"GET": ACCENT_BLUE, "POST": GREEN, "DELETE": RED}
for i, (method, path, desc, _) in enumerate(endpoints):
    col = i // 6
    row = i %  6
    y = 1.65 + row * 0.85
    x = 0.3 + col * 6.5
    mc = method_colors.get(method, MID_BLUE)
    add_rect(s, x, y, 0.9, 0.72, fill_color=mc)
    add_text(s, method, x+0.04, y+0.12, 0.82, 0.5,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x+0.95, y, 5.35, 0.72, fill_color=WHITE,
             line_color=mc, line_width=0.6)
    add_text(s, path, x+1.05, y+0.04, 5.1, 0.35, font_size=12, bold=True, color=mc)
    add_text(s, desc, x+1.05, y+0.38, 5.1, 0.3,  font_size=11, color=DARK_GRAY, wrap=True)

add_rect(s, 0.3, 6.76, 12.7, 0.45, fill_color=DARK_BLUE)
add_text(s, "👉  Interactive API docs available at:   http://localhost:8000/api/docs   (Swagger UI)",
         0.5, 6.79, 12.2, 0.38, font_size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — Frontend Pages
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Frontend Pages", "5 React pages built with Vite  ·  Runs on port 5173")

pages = [
    ("🏠", "Home\n/",             ["Hero section with project intro", "Feature cards with icons",
                                    "Live statistics from backend", "Start Analyzing CTA button"], ACCENT_BLUE),
    ("🔍", "Analyze\n/analyze",   ["Paste Text tab", "Upload File tab (TXT/PDF/DOCX)",
                                    "Loading spinner during analysis", "Error messages for validation"], MID_BLUE),
    ("📊", "Results\n/results/:id",["Similarity % with colour coding", "3 component score cards",
                                    "Top matching sources", "Download PDF Report button"], RGBColor(0x0E,0x73,0x48)),
    ("📋", "History\n/history",   ["List of all past analyses", "Date, score, classification",
                                    "Click to view full details", "Delete individual records"], RGBColor(0x71,0x3F,0x12)),
    ("📁", "Documents\n/documents",["View all reference docs", "Upload new document",
                                    "Word count & file info", "Delete reference docs"], RGBColor(0x4C,0x0D,0x99)),
]

for i, (icon, title, features, color) in enumerate(pages):
    x = 0.22 + i * 2.58
    add_rect(s, x, 1.65, 2.45, 0.72, fill_color=color)
    add_text(s, icon,  x+0.1, 1.68, 0.55, 0.65, font_size=24, color=WHITE)
    add_text(s, title, x+0.65, 1.72, 1.7,  0.62, font_size=12, bold=True, color=WHITE, wrap=True)
    for j, feat in enumerate(features):
        fy = 2.48 + j * 1.0
        add_rect(s, x, fy, 2.45, 0.88, fill_color=WHITE, line_color=color, line_width=0.6)
        add_text(s, f"• {feat}", x+0.1, fy+0.08, 2.28, 0.72,
                 font_size=11, color=DARK_GRAY, wrap=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — Folder Structure
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Project Folder Structure", "Clean modular architecture — separation of concerns")

backend_tree = [
    "backend/",
    "  app/",
    "    main.py          ← FastAPI entry point",
    "    config.py        ← Settings & thresholds",
    "    api/             ← Route handlers",
    "      routes_analysis.py",
    "      routes_documents.py",
    "      routes_reports.py",
    "    models/          ← DB models + schemas",
    "    services/        ← NLP algorithms",
    "      preprocessing.py",
    "      tfidf_similarity.py",
    "      ngram_similarity.py",
    "      fuzzy_matching.py",
    "      plagiarism_engine.py",
    "    utils/           ← Helpers",
    "  datasets/          ← Sample & test docs",
    "  tests/             ← 100+ pytest tests",
    "  requirements.txt",
    "  evaluation.py",
]
frontend_tree = [
    "frontend/",
    "  src/",
    "    main.jsx         ← Entry point",
    "    App.jsx          ← Routing setup",
    "    index.css        ← Global styles",
    "    pages/",
    "      Home.jsx/.css",
    "      Analyze.jsx/.css",
    "      Results.jsx/.css",
    "      History.jsx/.css",
    "      Documents.jsx/.css",
    "    components/",
    "      Navbar.jsx/.css",
    "    services/",
    "      api.js         ← All API calls",
    "  package.json",
    "  vite.config.js",
    "",
    "README.md",
    "DEMO_GUIDE.md",
    "PROJECT_REPORT.md",
]

add_rect(s, 0.3, 1.65, 6.2, 5.6, fill_color=DARK_GRAY)
add_text(s, "# Backend", 0.5, 1.68, 6.0, 0.4, font_size=14, bold=True, color=ACCENT_BLUE)
for i, line in enumerate(backend_tree):
    add_text(s, line, 0.45, 2.12 + i * 0.248, 6.0, 0.26,
             font_size=10, color=RGBColor(0xD1,0xFA,0xFF))

add_rect(s, 6.8, 1.65, 6.2, 5.6, fill_color=DARK_GRAY)
add_text(s, "# Frontend", 7.0, 1.68, 6.0, 0.4, font_size=14, bold=True, color=ACCENT_BLUE)
for i, line in enumerate(frontend_tree):
    add_text(s, line, 6.95, 2.12 + i * 0.248, 6.0, 0.26,
             font_size=10, color=RGBColor(0xD1,0xFA,0xFF))

# ══════════════════════════════════════════════════════════════════
# SLIDE 16 — How to Run
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "How to Run the Project", "Step-by-step startup guide")

steps_run = [
    ("STEP 1", "Start Backend Server",
     "cd backend\n.\\venv\\Scripts\\activate\nuvicorn app.main:app --reload",
     "Terminal 1  ·  Runs on http://localhost:8000", MID_BLUE),
    ("STEP 2", "Start Frontend Server",
     "cd frontend\nnpm run dev",
     "Terminal 2  ·  Runs on http://localhost:5173", RGBColor(0x0E,0x73,0x48)),
    ("STEP 3", "Upload Reference Documents",
     "Go to Documents page → Upload 5 files from:\nbackend/datasets/sample_documents/",
     "One-time setup  ·  Builds the reference database", RGBColor(0x71,0x3F,0x12)),
    ("STEP 4", "Run Plagiarism Analysis",
     "Go to Analyze → Paste text or upload file\nClick 'Analyze for Plagiarism'",
     "Results appear in seconds  ·  History saved automatically", RGBColor(0x4C,0x0D,0x99)),
]

for i, (step, title, cmd, note, color) in enumerate(steps_run):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 1.65 + row * 2.75
    add_rect(s, x, y, 0.85, 2.5, fill_color=color)
    add_text(s, step, x+0.04, y+0.9, 0.78, 0.9,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
    add_rect(s, x+0.88, y, 5.4, 2.5, fill_color=WHITE,
             line_color=color, line_width=0.8)
    add_text(s, title, x+1.0, y+0.08, 5.15, 0.45,
             font_size=15, bold=True, color=color)
    add_rect(s, x+0.9, y+0.55, 5.35, 1.3, fill_color=DARK_GRAY)
    add_text(s, cmd, x+1.05, y+0.62, 5.1, 1.12,
             font_size=11, color=RGBColor(0x86,0xEF,0xAC), wrap=True)
    add_text(s, f"ℹ  {note}", x+1.0, y+2.04, 5.15, 0.35,
             font_size=10, color=DARK_GRAY, italic=True, wrap=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 17 — Test Results & Evaluation
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Test Results & Evaluation", "4 test cases · 100% accuracy on all cases")

test_data = [
    ("high_similarity_text.txt",     "Directly copied content",   "87.57%", "Potential Plagiarism", GREEN, "PASS"),
    ("moderate_similarity_text.txt", "Paraphrased / reworded",    "37.13%", "Low Similarity",       GREEN, "PASS"),
    ("low_similarity_text.txt",      "Completely different topic", "3.34%",  "Mostly Original",      GREEN, "PASS"),
    ("original_text.txt",            "Original written content",  "14.96%", "Mostly Original",      GREEN, "PASS"),
]

headers = ["Test File", "Description", "Score", "Classification", "Result"]
col_x   = [0.35, 3.8, 7.3, 8.85, 11.7]
col_w   = [3.35, 3.4, 1.45, 2.75, 1.45]

add_rect(s, 0.3, 1.65, 12.7, 0.55, fill_color=DARK_BLUE)
for j, h in enumerate(headers):
    add_text(s, h, col_x[j], 1.68, col_w[j], 0.46,
             font_size=13, bold=True, color=WHITE)

for i, (fname, desc, score, classif, color, result) in enumerate(test_data):
    y = 2.28 + i * 0.82
    bg = WHITE if i % 2 == 0 else RGBColor(0xF0,0xF9,0xFF)
    add_rect(s, 0.3, y, 12.7, 0.75, fill_color=bg, line_color=ACCENT_BLUE, line_width=0.4)
    vals = [fname, desc, score, classif]
    for j, v in enumerate(vals):
        add_text(s, v, col_x[j], y+0.15, col_w[j], 0.5,
                 font_size=12, color=DARK_GRAY, wrap=True)
    add_rect(s, col_x[4], y+0.12, 1.2, 0.48, fill_color=GREEN)
    add_text(s, f"✅ {result}", col_x[4], y+0.17, 1.35, 0.42,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 0.3, 5.62, 12.7, 0.65, fill_color=GREEN)
add_text(s, "✅  Overall System Accuracy:   4 / 4 Tests PASSED   →   100%",
         0.5, 5.68, 12.2, 0.5, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

stats = [
    ("Avg TF-IDF", "49.51%",  ACCENT_BLUE),
    ("Avg N-Gram", "17.71%",  MID_BLUE),
    ("Avg Fuzzy",  "24.84%",  RGBColor(0x0E,0x73,0x48)),
    ("Accuracy",   "100%",    GREEN),
]
for i, (label, val, color) in enumerate(stats):
    x = 0.3 + i * 3.2
    add_rect(s, x, 6.4, 3.0, 0.8, fill_color=color)
    add_text(s, label, x+0.1, 6.44, 2.8, 0.34, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, val,   x+0.1, 6.76, 2.8, 0.35, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 18 — Limitations & Future Work
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Limitations & Future Enhancements", "Honest assessment + roadmap for improvement")

bullet_box(s,
    ["Local database only — no internet/web search",
     "English language support only",
     "Very short texts (<50 chars) may be inaccurate",
     "Cannot detect code or image plagiarism",
     "Heavy paraphrasing may reduce scores",
     "No user authentication system"],
    0.3, 1.65, 6.0, 5.6,
    font_size=14, title="⚠️  Current Limitations", title_color=RED)

future_short = [
    "User login system (teacher/student roles)",
    "Batch upload — analyze 50 files at once",
    "Email PDF reports automatically",
    "Dark mode UI theme",
]
future_long = [
    "BERT / sentence-transformers for better semantics",
    "Internet search integration",
    "Multi-language support (Hindi, Telugu…)",
    "LMS integration (Moodle, Google Classroom)",
    "Code plagiarism using AST comparison",
    "Blockchain certificates for originality",
]

bullet_box(s, future_short, 6.5, 1.65, 6.5, 2.55,
           font_size=13, title="🚀  Short-Term (1–3 months)", title_color=GREEN)
bullet_box(s, future_long,  6.5, 4.4, 6.5, 2.85,
           font_size=13, title="🔭  Long-Term (6+ months)", title_color=MID_BLUE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 19 — Viva Q&A
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_bg(s)
slide_header(s, "Key Viva Questions & Answers", "Top questions examiners ask — be prepared!")

qna = [
    ("Q1", "What is TF-IDF?",
     "Term Frequency × Inverse Document Frequency. Measures word importance in a document vs all documents. Rare but frequent words score highest."),
    ("Q2", "Why 3 algorithms?",
     "Each catches different plagiarism: TF-IDF=semantic, N-gram=exact phrases, Fuzzy=minor modifications. Combined = better accuracy."),
    ("Q3", "Why FastAPI over Flask?",
     "FastAPI is faster, has automatic Swagger docs, supports async, and has built-in Pydantic validation — better for modern APIs."),
    ("Q4", "What is cosine similarity?",
     "Measures angle between two TF-IDF vectors. cos(0°)=1.0 (identical). cos(90°)=0.0 (completely different)."),
    ("Q5", "How does it handle paraphrasing?",
     "TF-IDF captures topic-level similarity even when words change. If same vocab/concepts are used, TF-IDF still gives a high score."),
    ("Q6", "What are the limitations?",
     "No internet search, English only, no code/image detection, heavy paraphrasing may evade. Designed as assistant for human review."),
]

for i, (qnum, question, answer) in enumerate(qna):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 1.65 + row * 1.88
    add_rect(s, x, y, 0.75, 1.72, fill_color=ACCENT_BLUE)
    add_text(s, qnum, x+0.04, y+0.62, 0.68, 0.55,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x+0.78, y, 5.5, 1.72, fill_color=WHITE,
             line_color=ACCENT_BLUE, line_width=0.7)
    add_text(s, question, x+0.9, y+0.08, 5.25, 0.45,
             font_size=14, bold=True, color=MID_BLUE, wrap=True)
    add_text(s, answer, x+0.9, y+0.56, 5.25, 1.1,
             font_size=12, color=DARK_GRAY, wrap=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 20 — Conclusion
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=DARK_BLUE)
add_rect(s, 0, 2.5, 13.33, 0.08, fill_color=ACCENT_BLUE)
add_rect(s, 0, 2.6, 13.33, 0.04, fill_color=RGBColor(0x93,0xC5,0xFD))

add_text(s, "Conclusion", 0.6, 0.25, 12.0, 0.9,
         font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "What we built · What we learned · What comes next",
         0.6, 1.1, 12.0, 0.5, font_size=18,
         color=RGBColor(0xBA,0xD4,0xFB), align=PP_ALIGN.CENTER)

achievements = [
    "✅  3-algorithm NLP pipeline (TF-IDF + N-gram + Fuzzy Matching)",
    "✅  Full-stack web app — FastAPI backend + React.js frontend",
    "✅  100% accuracy on 4 evaluation test cases",
    "✅  PDF report generation with colour-coded analysis",
    "✅  SQLite database with full analysis history",
    "✅  100+ unit tests for reliability",
    "✅  Supports TXT, PDF, and DOCX file formats",
]
add_rect(s, 0.5, 2.8, 12.3, 3.3, fill_color=RGBColor(0x0D,0x1A,0x4A))
add_text(s, "Project Achievements", 0.7, 2.88, 11.8, 0.45,
         font_size=16, bold=True, color=ACCENT_BLUE)
for i, ach in enumerate(achievements):
    add_text(s, ach, 0.8, 3.32 + i * 0.37, 11.8, 0.34,
             font_size=13, color=WHITE, wrap=True)

add_rect(s, 0.5, 6.22, 12.3, 0.65, fill_color=RGBColor(0x0F,0x0F,0x4A))
add_text(s,
    "This project demonstrates how multiple AI/NLP techniques can be combined to solve "
    "a real-world problem — detecting academic plagiarism automatically.",
    0.7, 6.27, 11.8, 0.55, font_size=13,
    color=RGBColor(0x7D,0xD3,0xFC), wrap=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
output_path = r"C:\Users\bhanu\OneDrive\Desktop\Combat\AI_Plagiarism_Detector_Presentation.pptx"
prs.save(output_path)
print(f"✅  Presentation saved to:\n    {output_path}")
print(f"    Slides: {len(prs.slides)}")
